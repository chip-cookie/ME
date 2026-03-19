"""
Learning Module Service
"""
import json
import logging
from pathlib import Path
from typing import List, Optional
from PIL import Image
import pytesseract
from sqlalchemy.orm import Session

from app.core.config import get_settings
from app.modules.learning.models import LearningLog, SuccessfulExample
from app.modules.style.models import StyleProfile  # Cross-module dependency

settings = get_settings()
logger = logging.getLogger(__name__)


class LearningService:
    def __init__(self, db: Session):
        self.db = db
        self.learning_dir = Path(settings.learning_dir)
        self._ensure_dirs()
    
    def _ensure_dirs(self):
        dirs = [
            self.learning_dir / "images" / "raw",
            self.learning_dir / "images" / "processed",
            self.learning_dir / "styles",
            self.learning_dir / "styles",
            self.learning_dir / "history",
            Path(settings.successful_examples_dir),
        ]
        for d in dirs:
            d.mkdir(parents=True, exist_ok=True)
    
    async def learn_from_images(
        self, 
        images: List[bytes], 
        style_name: str,
        meta_tags: Optional[dict] = None
    ) -> dict:
        extracted_texts = []
        
        for i, img_bytes in enumerate(images):
            text = self._extract_text_from_image(img_bytes)
            if text:
                extracted_texts.append(text)
                processed_path = self.learning_dir / "images" / "processed" / f"{style_name}_{i}.json"
                with open(processed_path, "w", encoding="utf-8") as f:
                    json.dump({"text": text, "index": i}, f, ensure_ascii=False, indent=2)
        
        if not extracted_texts:
            return {
                "style_id": None,
                "extracted_features": {},
                "status": "failed",
                "message": "이미지에서 텍스트를 추출할 수 없습니다."
            }
        
        features = self._extract_style_features(extracted_texts)
        if meta_tags:
            features["meta"] = meta_tags
        
        style = self._create_or_update_style(style_name, features)
        self._log_learning(style.id, len(images))
        
        return {
            "style_id": style.id,
            "extracted_features": features,
            "status": "success",
            "message": f"스타일 '{style_name}'이 성공적으로 학습되었습니다."
        }
    
    async def ingest_from_local(self) -> dict:
        import shutil
        import os
        from pypdf import PdfReader
        from docx import Document
        
        import_dir = self.learning_dir / "import"
        if not import_dir.exists():
            import_dir.mkdir(parents=True)
            return {"status": "ok", "message": "Import directory created. Please put files in 'data/learning/import/{style_name}'"}
        
        raw_dir = self.learning_dir / "images" / "raw" # Using existing raw dir structure for now
        processed_dir = self.learning_dir / "images" / "processed"
        
        results = {}
        
        # Iterate over subdirectories (styles)
        for style_path in import_dir.iterdir():
            if not style_path.is_dir():
                continue
                
            style_name = style_path.name
            good_texts = []
            bad_texts = []
            
            # Helper to process a directory
            def process_dir(directory: Path, target_list: List[str]):
                if not directory.exists(): return
                for file_path in directory.iterdir():
                    if file_path.name.startswith('.'): continue
                    text = self._extract_text_from_file(file_path)
                    if text:
                        target_list.append(text)

            # Check for subdirectories 'good' and 'bad'
            good_dir = style_path / "good"
            bad_dir = style_path / "bad"
            
            if good_dir.exists() or bad_dir.exists():
                process_dir(good_dir, good_texts)
                process_dir(bad_dir, bad_texts)
            else:
                # If no subdirs, treat root files as 'good' examples (default behavior)
                process_dir(style_path, good_texts)
            
            # Combine for feature extraction (using only good texts for style analysis)
            if good_texts:
                features = self._extract_style_features(good_texts)
                features["good_examples"] = good_texts[:5] # Store top 5
                features["bad_examples"] = bad_texts[:5]   # Store top 5
                
                style = self._create_or_update_style(style_name, features)
                self._log_learning(style.id, len(good_texts) + len(bad_texts))
                results[style_name] = len(good_texts) + len(bad_texts)
        
        return {
            "status": "success",
            "results": results,
            "message": f"Successfully ingested styles: {', '.join(results.keys())}"
        }

    async def auto_ingest_and_classify(self, style_name: str, criteria: str) -> dict:
        """AI가 기준에 따라 자동으로 Good/Bad 분류 후 학습"""
        from ollama import AsyncClient
        
        import_dir = self.learning_dir / "import" / style_name
        if not import_dir.exists():
            # Check interview directory
            interview_dir = self.learning_dir / "interview" / style_name
            if interview_dir.exists():
                import_dir = interview_dir
            else:
                import_dir.mkdir(parents=True)
                return {"status": "error", "message": f"폴더가 없습니다. 파일을 'data/learning/import/{style_name}' 또는 'data/learning/interview/{style_name}'에 넣어주세요."}
        
        # Extract text from all files (공통 메서드 사용)
        file_texts = []
        for file_path in import_dir.iterdir():
            if file_path.is_dir() or file_path.name.startswith('.'): continue
            text = self._extract_text_from_file(file_path)
            if text:
                file_texts.append({"file": file_path.name, "text": text})
        
        if not file_texts:
            return {"status": "error", "message": "추출 가능한 텍스트가 없습니다."}
        
        # AI-based classification
        good_texts = []
        bad_texts = []
        
        client = AsyncClient(host=settings.ollama_base_url)
        
        for item in file_texts:
            prompt = f"""다음 글이 아래 기준에 얼마나 부합하는지 0~100 점수로만 답해줘. 숫자만 출력해.
기준: {criteria}
글: {item['text'][:500]}"""
            try:
                response = await client.chat(
                    model=settings.ollama_model,
                    messages=[{'role': 'user', 'content': prompt}]
                )
                score_text = response.message.content.strip()
                # Extract number from response
                import re
                match = re.search(r'\d+', score_text)
                score = int(match.group()) if match else 50
                
                if score >= 50:
                    good_texts.append(item['text'])
                else:
                    bad_texts.append(item['text'])
                logger.info(f"[{item['file']}] Score: {score} -> {'Good' if score >= 50 else 'Bad'}")
            except Exception as e:
                logger.error(f"AI scoring failed for {item['file']}: {e}")
                good_texts.append(item['text'])  # Default to good on error
        
        # Create style with classified examples
        if good_texts or bad_texts:
            features = self._extract_style_features(good_texts if good_texts else bad_texts)
            features["good_examples"] = good_texts[:5]
            features["bad_examples"] = bad_texts[:5]
            
            style = self._create_or_update_style(style_name, features)
            self._log_learning(style.id, len(good_texts) + len(bad_texts))
            
            return {
                "status": "success",
                "style_id": style.id,
                "good_count": len(good_texts),
                "bad_count": len(bad_texts),
                "message": f"자동 분류 완료: Good {len(good_texts)}개, Bad {len(bad_texts)}개"
            }
        
        return {"status": "error", "message": "분류할 텍스트가 없습니다."}

    def _extract_text_from_image(self, img_bytes: bytes) -> str:
        try:
            from io import BytesIO
            img = Image.open(BytesIO(img_bytes))
            text = pytesseract.image_to_string(img, lang='kor+eng')
            return text.strip()
        except Exception:
            return ""

    def _extract_text_from_file(self, file_path: Path) -> str:
        """파일에서 텍스트를 추출합니다. 지원 포맷: jpg, png, pdf, docx, txt, pptx, hwp"""
        from pypdf import PdfReader
        from docx import Document
        
        text = ""
        suffix = file_path.suffix.lower()
        
        try:
            if suffix in ['.jpg', '.jpeg', '.png']:
                with open(file_path, "rb") as f:
                    text = self._extract_text_from_image(f.read())
            elif suffix == '.pdf':
                reader = PdfReader(str(file_path))
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            elif suffix == '.docx':
                doc = Document(str(file_path))
                text = "\n".join([para.text for para in doc.paragraphs])
            elif suffix in ['.txt', '.md']:
                with open(file_path, "r", encoding="utf-8") as f:
                    text = f.read()
            elif suffix == '.pptx':
                from pptx import Presentation
                prs = Presentation(str(file_path))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif suffix == '.hwp':
                import subprocess
                result = subprocess.run(
                    [settings.hwp5txt_path, str(file_path)],
                    capture_output=True, text=True, timeout=30
                )
                if result.returncode == 0:
                    text = result.stdout
        except Exception as e:
            logger.error(f"파일 처리 오류 {file_path}: {e}")
        
        return text.strip()
    
    def _extract_style_features(self, texts: List[str]) -> dict:
        all_text = "\n".join(texts)
        sentences = [s.strip() for s in all_text.split('.') if s.strip()]
        words = all_text.split()
        
        return {
            "tone": self._analyze_tone(all_text),
            "structure": self._analyze_structure(texts),
            "expressions": self._extract_expressions(all_text),
            "stats": {
                "avg_sentence_length": len(words) / max(len(sentences), 1),
                "total_sentences": len(sentences),
                "sample_count": len(texts)
            }
        }
    
    def _analyze_tone(self, text: str) -> dict:
        return {
            "formality": "formal" if any(w in text for w in ["습니다", "입니다"]) else "casual",
            "politeness": "polite" if any(w in text for w in ["드립니다", "감사"]) else "neutral",
        }
    
    def _analyze_structure(self, texts: List[str]) -> dict:
        return {
            "uses_bullet_points": any("•" in t or "-" in t for t in texts),
            "uses_numbering": any("1." in t for t in texts),
            "paragraph_style": "short" if all(len(t) < 500 for t in texts) else "long"
        }
    
    def _extract_expressions(self, text: str) -> dict:
        endings = [e for e in ["습니다", "입니다", "합니다", "요"] if e in text]
        return {"common_endings": endings[:5]}
    
    def _create_or_update_style(self, name: str, features: dict) -> StyleProfile:
        style = self.db.query(StyleProfile).filter(StyleProfile.name == name).first()
        if style:
            style.tone_patterns = features.get("tone", {})
            style.structure_rules = features.get("structure", {})
            style.expression_dict = features.get("expressions", {})
            style.good_examples = features.get("good_examples", [])
            style.bad_examples = features.get("bad_examples", [])
            style.sample_count += features.get("stats", {}).get("sample_count", 0)
        else:
            style = StyleProfile(
                name=name,
                description=f"로컬 파일에서 학습된 스타일",
                tone_patterns=features.get("tone", {}),
                structure_rules=features.get("structure", {}),
                expression_dict=features.get("expressions", {}),
                good_examples=features.get("good_examples", []),
                bad_examples=features.get("bad_examples", []),
                sample_count=features.get("stats", {}).get("sample_count", 0)
            )
            self.db.add(style)
        self.db.commit()
        self.db.refresh(style)
        return style
    
    def _log_learning(self, style_id: int, samples_count: int):
        log = LearningLog(samples_count=samples_count, metrics={"style_id": style_id})
        self.db.add(log)
        self.db.commit()

    async def ingest_successful_examples(self) -> dict:
        """성공적인 자소서 예시를 분석하고 저장합니다."""
        from ollama import AsyncClient
        
        target_dir = Path(settings.successful_examples_dir)
        if not target_dir.exists():
            target_dir.mkdir(parents=True)
            return {"status": "error", "message": f"폴더가 없습니다. 파일을 '{target_dir}'에 넣어주세요."}
            
        processed_count = 0
        client = AsyncClient(host=settings.ollama_base_url)
        
        for file_path in target_dir.iterdir():
            if file_path.is_dir() or file_path.name.startswith('.'): continue
            
            # Check if likely text
            if file_path.suffix.lower() not in ['.txt', '.md']:
                # For now only support simple text files for stability, or expand if needed
                # If users want other formats, they can use the general ingest which supports more
                # but for this specific 'success' folder let's support txt/md for now
                pass

            try:
                content = ""
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                
                if not content.strip(): continue
                
                # Check duplication
                existing = self.db.query(SuccessfulExample).filter(SuccessfulExample.content == content).first()
                if existing:
                    continue
                
                # Analyze with Ollama
                prompt = f"""다음은 합격한 자기소개서 예시입니다. 이 글이 왜 합격했는지, 어떤 점이 훌륭한지 분석해주세요.
특히 다음 내용을 포함해서 JSON 형식으로 답해주세요:
1. structure: 글의 구조적 특징
2. flow: 논리의 흐름
3. key_expressions: 주요 표현이나 어휘
4. summary: 전체적인 총평

자소서:
{content}
"""
                response = await client.chat(
                    model=settings.ollama_model,
                    messages=[{'role': 'user', 'content': prompt}],
                    format="json"
                )
                
                analysis = json.loads(response.message.content)
                
                # Store
                example = SuccessfulExample(
                    content=content,
                    analysis=analysis,
                    category="general" # Can be refined later
                )
                self.db.add(example)
                processed_count += 1
                
            except Exception as e:
                logger.error(f"Error processing {file_path}: {e}")
        
        self.db.commit()
        return {
            "status": "success",
            "processed_count": processed_count,
            "message": f"{processed_count}개의 합격 자소서를 분석하여 저장했습니다."
        }

